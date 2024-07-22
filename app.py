
import nltk
nltk.download('stopwords')
nltk.download('punkt')
import os
import io
from flask import Flask, render_template_string, request, session
from flask_sqlalchemy import SQLAlchemy
# from pyngrok import ngrok
from transformers import pipeline, BartTokenizer, BartForConditionalGeneration
from PIL import Image
import fitz  # PyMuPDF
import pdfplumber
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import torch
import torchvision.transforms as transforms
from torchvision import models
#import pytesseract
import groq
from spire.doc import *
from spire.doc.common import *
import re
from docx import Document
from datetime import datetime
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
import PIL.Image
from datetime import datetime
# Step 3: Define the Flask app
app = Flask(__name__)
app.secret_key = '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'  # Required for session management

# SQLAlchemy configuration
app.config['SECRET_KEY'] = 'your_secret_key'
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///summaries.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db = SQLAlchemy(app)

# Define the Summary model
class Summary(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    content = db.Column(db.Text, nullable=False)
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)

# Create the database and tables
with app.app_context():
    db.create_all()
class QuestionAnswer(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    question = db.Column(db.String(500), nullable=False)
    answer = db.Column(db.Text, nullable=False)
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)

    def __repr__(self):
        return f'<QuestionAnswer {self.id}>'

@app.before_first_request
def create_tables():
    db.create_all()

# Ensure the uploads directory exists
os.makedirs("uploads", exist_ok=True)

# Define the labels for graph types
graph_labels = ["bar", "line", "pie", "scatter"]

# Load and configure ResNet50 model for graph classification
graph_model = models.resnet50(pretrained=True)
num_classes = len(graph_labels)
graph_model.fc = torch.nn.Linear(graph_model.fc.in_features, num_classes)
graph_model.eval()

# def summarize_text_from_word(text):
#     chunks = [text[i:i + 1000] for i in range(0, len(text), 1000)]
#     return "".join(summarizer(chunk, max_length=130, min_length=40, do_sample=False)[0]["summary_text"] for chunk in chunks)
def preprocess_text(text):
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces/newlines with a single space
    text = text.strip()  # Remove leading and trailing whitespace
    stop_words = set(stopwords.words('english'))
    word_tokens = word_tokenize(text)
    filtered_text = [word for word in word_tokens if word.lower() not in stop_words]
    text = ' '.join(filtered_text)
    return text

# Function to preprocess the image
def preprocess_image(image_path):
    input_image = Image.open(image_path).convert("RGB")
    preprocess = transforms.Compose([
        transforms.Resize(256),
        transforms.CenterCrop(224),
        transforms.ToTensor(),
        transforms.Normalize(mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225]),
    ])
    input_tensor = preprocess(input_image)
    return input_tensor.unsqueeze(0)

# Function to extract and save images from PPT
def extract_images_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    image_files = []
    for slide_num, slide in enumerate(prs.slides):
        for shape_num, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = io.BytesIO(image.blob)
                image_file_name = f"slide_{slide_num+1}_image_{shape_num+1}.{image.ext}"
                with open(image_file_name, "wb") as f:
                    f.write(image_bytes.read())
                image_files.append(image_file_name)
    return image_files

# Function to extract text from an image using OCR
#def extract_text_from_image(image_path):
 #   image = PIL.Image.open(image_path)
 #   return pytesseract.image_to_string(image)

# Extract text from PDF
def extract_text_from_pdf(pdf_path):
    with pdfplumber.open(pdf_path) as pdf:
        return "\n".join(page.extract_text() for page in pdf.pages)

# Extract text from PPT
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    return " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

def read_word_file(file_path):
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

# Summarize text
# def summarize_text(text):
#     chunks = [text[i:i + 1000] for i in range(0, len(text), 1000)]
#     return "".join(summarizer(chunk, max_length=102, min_length=40, do_sample=False)[0]["summary_text"] for chunk in chunks)

# def summarize_text_from_ppt(text):
#     chunks = [text[i:i + 1000] for i in range(0, len(text), 1000)]
#     return "".join(summarizer(chunk, max_length=130, min_length=30, do_sample=False)[0]["summary_text"] for chunk in chunks)

# Extract limited images from PDF
def extract_limited_images_from_pdf(pdf_path, image_dir, limit=4):
    if not os.path.exists(image_dir):
        os.makedirs(image_dir)

    doc = fitz.open(pdf_path)
    image_captions = []
    image_count = 0

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            if image_count >= limit:
                return image_captions
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_path = os.path.join(image_dir, f"image{page_num}_{img_index}.png")
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            image_captions.append(f"Image from page {page_num + 1}, image index {img_index + 1}")
            image_count += 1
    return image_captions

# Merge captions and text
def merge_captions_and_text(text, image_captions):
    return "\n\n".join([text] + image_captions)

def chunk_text(text, chunk_size=1000):
    """Splits the text into smaller chunks."""
    for i in range(0, len(text), chunk_size):
        yield text[i:i + chunk_size]

def summ(text, chunk_size=1000):
    summaries = []
    api_key = "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"
    client = groq.Client(api_key=api_key)

    for chunk in chunk_text(text, 1000):
        response = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": "Please summarize the following text:\n\n" + chunk,
                }
            ],
            model="llama3-8b-8192",
        )
        summaries.append(response.choices[0].message.content)

    # Combine the summaries of all chunks
    combined_summary = " ".join(summaries)
    return combined_summary

def format_bullet_points(text):
    formatted_response = text.replace('**', '')  # Add a newline before each numbered bullet point
    formatted_text = re.sub(r'\.\s*(\d+\.)', r'.\n\n\1', formatted_response)
    return formatted_text

# Answer question
def answer_question(question, context):
    api_key = "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"  # Replace with your Groq API key
    client = groq.Client(api_key=api_key)
    response = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": context + "\n\nQuestion: " + question,
            }
        ],
        model="llama3-8b-8192",
    )
    return response.choices[0].message.content


html_template = """
   <!DOCTYPE html>
   <html lang="en">
   <head>
       <meta charset="UTF-8">
       <meta name="viewport" content="width=device-width, initial-scale=1.0">
       <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.0-alpha3/dist/css/bootstrap.min.css" rel="stylesheet" integrity="sha384-KK94CHFLLe+nY2dmCWGMq91rCGa5gtU4mk92HdvYe+M/SXH301p5ILy+dN9+nJOZ" crossorigin="anonymous">
       <script src="https://cdn.jsdelivr.net/npm/bootstrap@5.3.1/dist/js/bootstrap.bundle.min.js" integrity="sha384-HwwvtgBNo3bZJJLYd8oVXjrBZt8cqVSpeBNS5n7C8IVInixGAoxmnlMuBnhbgrkm" crossorigin="anonymous"></script>
       <title>Document Processor</title>
       <style>
           body {
               /* background-image: #ffff; */
               font-family: sans-serif;
               background-color: rgb(4,4,102);
               height:100%;
               width:100%;

           }
           .main {
               background-image: linear-gradient(to top, #cfd9df 0%, #e2ebf0 100%);
               /* max-width: 620px; */
               /* background-color: white; */
               margin: 40px;
               padding: 20px;
               border: 1px solid #ddd;
               border-radius: 5px;
               box-shadow: 0 0 10px rgba(0, 0, 0, 0.1);
               color: #533a8c;
               height: 645px;
           }
           .header {
               display: flex;
               justify-content: space-between;
               margin-bottom: 10px;
               color: rgb(4,4,102);
           }
           .model {
               font-weight: bold;
           }
           .chat-window {
               border: 1px solid rgb(4,4,102);
               padding: 10px;
               height: 79%;
               overflow-y: scroll;
               margin-bottom: 10px;
           }
           .chat-message {
               margin-bottom: 10px;
           }
           .chat-message p {
               margin: 0;
               justify-content:flex-end;
           }
          .upload-file{
            display:flex;
            flex-direction:row;
            justify-content:space-between;
            align-items:center;
            margin-top:0px;
          }
           .user-query {

            padding: 10px;
            border-radius: 10px;
            margin-bottom: 10px;
            text-align: center;
            display:flex;
            margin-right:300px;
            width:100%;
            max-width:600px;
            justify-content:flex-end;
            align-items:center;
           }
           .chat-response {
               text-align: left;
               background-color:#8f8f8f
           }
           .output-status {
               font-size: 0.9em;
               color: gray;
               text-align: center;
               margin-top: 10px;
           }
           .input-area {
               display: flex;
               flex-direction: row;
              #  align-items: center;
              #  justify-content: flex-end;
               width: 100%;
              #  align-content: center;
           }
           .input-area input[type="text"] {
               padding: 10px;
               border: 1px solid #ddd;
               border-radius: 5px;
              #  width: 100%;
              #  max-width:600px;
              #  margin-left: 0px;
               justify-content: center;
               align-items: center;
               align-content: center;
               flex-grow: 1;
              margin-right: 50px;

           }

           .input-area button[type="upload"] {
               border: 1px solid #ddd;
               border-radius: 15px;
               background-color: #007bff;
               color: white;
               cursor: pointer;
               width:100px;
              height:40px;
           }
           .input-area button {
               padding: 10px;
               border: 1px solid #ddd;
               border-radius: 5px;
               background-color: #007bff;
               color: white;
               cursor: pointer;
           }
           .input-area button:hover {
               background-color: linear-gradient(to right, #ff8177 0%, #ff867a 0%, #ff8c7f 21%, #f99185 52%, #cf556c 78%, #b12a5b 100%);
           }
           .input-area input[type='file'] {
              display:none
               font-size: large;
           }
           .btn_submit{
            # justify-content-right;
            width:90px;
            height:42px;
            border-radius:15px 15px 15px 15px;
            # margin-right:30px
           }

             /* Responsive design */
        @media (min-width: 600px) {
            .input-area {
                flex-direction: row;
                justify-content: space-between;
                width: 100%;
            }
            .input-area input[type="text"] {
              margin-top:0px;
                width: 600px; /* 70% width on medium and large screens */
                max-width: 600px;
            }
        #status-message {
            display: flex;
            align-items: center;
        }
        #status-message .dots {
            display: inline-block;
            margin-left: 5px;
            width: 5px;
            height: 5px;
            border-radius: 50%;
            background-color: #007bff;
            animation: blink 1s infinite both;
        }
        #status-message .dots:nth-child(2) {
            animation-delay: 0.2s;
        }
        #status-message .dots:nth-child(3) {
            animation-delay: 0.4s;
        }
        @keyframes blink {
            0%, 80%, 100% {
                opacity: 0;
            }
            40% {
                opacity: 1;
            }
        }
        #boxing{
          background-color:#0f0f0f
          color:white
        }
       </style>
           <script>
        function showProcessingMessage() {
            document.getElementById('status-message').innerHTML = 'Summarizing your document<span class="dots">.</span><span class="dots">.</span><span class="dots">.</span>';
        }
    </script>
   </head>
   <body>
       <div class="main">
           <div class="header">
               <span class="model"><h2>Model: ChatG</h2></span>
               <span><h3><i>So what's on your mind right now?</i></h3></span>
               <!-- <span><h3><i>SO WHAT'S ON YOUR MIND NOW?</i></h3></span> -->

           </div>
           <div class="chat-window" id="chat-window">
<p id="status-message"></p>
               {% if query %}
               <div class="chat-message user-query" id="boxing">
               <p style="background-color:rgb(4,4,102);color:#ffff;border-radius:20px 20px 20px 20px;padding:15px"><strong>Question: {{ query }}</strong> </p>
               </div>
               {% endif %}
                {% if result_sum %}
                       <p style="background-color:rgb(4,4,102);color:#ffff;border-radius:20px 20px 20px 20px;padding:15px"><b>{{ result_sum }}</b></p><br><br>
                    {% endif %}
               {% if result %}
               <div class="chat-message user-response">
                   <p style="background-color:rgb(4,4,102);color:#ffff;border-radius:20px 20px 20px 20px;padding:15px"><strong>Response-> {{ result }}</strong> </p>
               </div>
               {% endif %}
               <div class="output-status">Upload the PDF,PPT OR Word and get its analysis and ask your questions{{ status }}</div>
           </div>
           <div style="display: flex; align-items: center;">
           <div class="input-area">
  <div style="display:flex;flex-direction:row;justify-content:flex-center;align-items:center;margin-left:20%">
    <div class="upload-file">
      <form method="post" enctype="multipart/form-data">
        <input type="file" name="file" id="file" accept=".pdf,.ppt,.pptx,.doc,.docx" hidden>
        <label for="file">
          <svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24" fill="currentColor" width="24" height="24"><path d="M19 13h-6v6h-2v-6h-6v-2h6v-6h2v6h6v2zm-8-9c-1.1 0-2.9-2 2s.9 2 2 2 2-.9 2-2-.9-2-2-2zm0 12c-1.1 0-2.9-2 2s.9 2 2 2 2-.9 2-2-.9-2-2-2z"/></svg>
        </label>
        <button type="submit" onclick="showProcessingMessage()">Upload</button>
      </form>
    </div>
    <form method="post" enctype="multipart/form-data">
    <input style="border-radius:35px;padding:20px;margin-left:20px;margin-right:20px" type="text" name="user_query" placeholder="Ask a question...">
    <button class="btn_submit" type="submit">Ask</button>
    </form>
  </div>
</div>
           </div>
       </div>
   </body>
   </html>
   """
@app.route("/", methods=["GET", "POST"])
def home():
    status = ".Leverage the chatbot's analysis tool for your documents.\nUpload your Document to access chat."
    result = None
    query = None
    result_sum = None
    history = QuestionAnswer.query.order_by(QuestionAnswer.timestamp.desc()).all()

    if request.method == "POST":
        if "file" in request.files:
            file = request.files["file"]
            if file.filename:
                # Save the uploaded file
                file_path = os.path.join("uploads", file.filename)
                file.save(file_path)

                # Determine the file type
                if file.filename.endswith(".pdf"):
                    print("uploaded pdf n now extracting text")
                    text = extract_text_from_pdf(file_path)
                    print("text extracted n now summary time")
                    summary = summ(text)
                    summary = format_bullet_points(summary)
                    new_summary = Summary(content=summary)
                    db.session.add(new_summary)
                    db.session.commit()

                    session['summary_id'] = new_summary.id
                    result_sum = summary
                    status = "PDF processed and summary generated."
                    print("summary:- ",result_sum)

                elif file.filename.endswith((".ppt", ".pptx")):
                       print("uploaded ppt n now extracting text")
                       text = extract_text_from_ppt(file_path)
                       print("text extracted n now summary time")
                       summary = summ(text)
                       print("summary of text done n extracting images now")
                       extracted_images = extract_images_from_ppt(file_path)
                       descriptions = [summary]
                       print("extracting text from images")
                       #for image_file in extracted_images:
                        #   extracted_text = extract_text_from_image(image_file)
                        #   descriptions.append(extracted_text)
                       result = "\n".join(descriptions)
                       # Save summary to the database
                       summary = format_bullet_points(summary)
                       result_sum = Summary(content=summary)
                       db.session.add(result_sum)
                       db.session.commit()
                       session['summary_id'] = result_sum.id
                       status = "PPT processed and descriptions generated."

                elif file.filename.endswith((".doc", ".docx")):
                       print("uploaded word doc n extracting text now")
                       text = read_word_file(file_path)
                       print("text done now summary time")
                       summary = summ(text)

                       # Save summary to the database
                       new_summary = Summary(content=summary)
                       db.session.add(new_summary)
                       db.session.commit()

                       session['summary_id'] = new_summary.id
                       summary = format_bullet_points(summary)
                       result_sum = summary
                       status = "Word document processed and summary generated."

        if "user_query" in request.form:
            user_query = request.form.get("user_query")
            summary_id = session.get('summary_id', None)
            if summary_id:
                summary = Summary.query.get(summary_id).content

            if user_query:
                if user_query.lower() == 'question history':
                    query = user_query
                    qa_entries = QuestionAnswer.query.order_by(QuestionAnswer.timestamp.desc()).all()
                    result = "\n".join([f"Q: {qa.question}\nA: {qa.answer}" for qa in qa_entries])
                elif user_query.lower() == 'summary history':
                    query = user_query
                    summaries = Summary.query.order_by(Summary.id.desc()).all()
                    result = "\n".join([f"Summary {i+1}:\n{summary.content}" for i, summary in enumerate(summaries)])
                elif user_query.startswith('[integrate]'):
                    query = user_query
                    summaries = Summary.query.order_by(Summary.timestamp).all()
                    combined_summary = " ".join(summary.content for summary in summaries)
                    combined_summary = preprocess_text(combined_summary)
                    user_query = user_query[len('[integrate]'):].strip()  # Remove the [integrate] keyword
                    answer = answer_question(user_query, combined_summary)
                    result = answer if answer else "Sorry, I couldn't find an answer."
                else:
                    query = user_query
                    answer = answer_question(user_query, summary)
                    result = answer if answer else "Sorry, I couldn't find an answer."

                result = format_bullet_points(result)
                qa_entry = QuestionAnswer(question=user_query, answer=result)
                db.session.add(qa_entry)
                db.session.commit()
            else:
                result = "Please upload a file and generate a summary first."




    return render_template_string(html_template, result=result, query=query, result_sum=result_sum, status=status, history=history)


# Step 9: Run the Flask app
if __name__ == "__main__":
    # Run the Flask application
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 5000)))
